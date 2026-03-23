"""
MindPulse PPT Generator — Professional version
Uses template's visual style (teal #277884, Arial, 4:3) with full design control.
No placeholder conflicts — all shapes are explicitly positioned.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

prs = Presentation("D:/another-project/mini.pptx")

# Remove old template slides
while len(prs.slides) > 0:
    rId = prs.slides._sldIdLst[0].get(qn("r:id"))
    prs.part.drop_rel(rId)
    prs.slides._sldIdLst.remove(prs.slides._sldIdLst[0])

# Colors matching template
TEAL = RGBColor(0x27, 0x78, 0x84)
DARK = RGBColor(0x2C, 0x2C, 0x2C)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT = RGBColor(0xF0, 0xF0, 0xF0)
RED = RGBColor(0xCC, 0x33, 0x33)
MID_GRAY = RGBColor(0x66, 0x66, 0x66)
BG_TEAL = RGBColor(0x1A, 0x3A, 0x4A)

SW = prs.slide_width  # 10 inches
SH = prs.slide_height  # 7.5 inches


def _clear_slide(slide):
    """Remove all existing shapes from a slide."""
    spTree = slide.shapes._spTree
    # Keep the spTree itself but remove shape children
    for sp in list(spTree):
        if sp.tag.endswith("}sp") or sp.tag.endswith("}grpSp"):
            spTree.remove(sp)


def _add_rect(slide, left, top, w, h, fill=None, line=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shape.shadow.inherit = False
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    return shape


def _add_rounded(slide, left, top, w, h, fill=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.shadow.inherit = False
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    return shape


def _set_para(run, text, size=Pt(12), bold=False, color=DARK, font="Arial"):
    run.text = text
    run.font.name = font
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color


def _add_text_box(slide, left, top, w, h, lines, default_size=Pt(12)):
    """Add a text box with multiple lines. Each line = (text, size, bold, color)."""
    txBox = slide.shapes.add_textbox(left, top, w, h)
    txBox.shadow.inherit = False
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)

    for i, line_data in enumerate(lines):
        if isinstance(line_data, str):
            text, sz, bld, clr = line_data, default_size, False, DARK
        else:
            text = line_data[0]
            sz = line_data[1] if len(line_data) > 1 else default_size
            bld = line_data[2] if len(line_data) > 2 else False
            clr = line_data[3] if len(line_data) > 3 else DARK

        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(2)
        run = p.add_run()
        _set_para(run, text, size=sz, bold=bld, color=clr)
    return txBox


def _add_card(slide, left, top, w, h, lines, bg=WHITE, border_color=None):
    """Add a card (rounded rect with text inside)."""
    card = _add_rounded(slide, left, top, w, h, fill=bg)
    if border_color:
        card.line.color.rgb = border_color
        card.line.width = Pt(1.5)
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.1)
    tf.margin_bottom = Inches(0.1)
    for i, line_data in enumerate(lines):
        if isinstance(line_data, str):
            text, sz, bld, clr = line_data, Pt(11), False, DARK
        else:
            text = line_data[0]
            sz = line_data[1] if len(line_data) > 1 else Pt(11)
            bld = line_data[2] if len(line_data) > 2 else False
            clr = line_data[3] if len(line_data) > 3 else DARK
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(1)
        p.space_after = Pt(1)
        run = p.add_run()
        _set_para(run, text, size=sz, bold=bld, color=clr)
    return card


def _add_table(slide, left, top, w, h, headers, rows, col_widths=None):
    """Add a styled table."""
    tbl_shape = slide.shapes.add_table(len(rows) + 1, len(headers), left, top, w, h)
    tbl = tbl_shape.table

    # Header row
    for i, hdr in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.text = hdr
        cell.fill.solid()
        cell.fill.fore_color.rgb = TEAL
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                _set_para(run, hdr, size=Pt(10), bold=True, color=WHITE)

    # Data rows
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = tbl.cell(r + 1, c)
            cell.text = val
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r % 2 == 0 else LIGHT
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    _set_para(run, val, size=Pt(9), color=DARK)

    # Column widths
    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = cw

    return tbl_shape


def _set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


# ================================================================
# SLIDE 1: TITLE
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[1])  # BLANK
_clear_slide(s)

# Full teal background
_add_rect(s, Inches(0), Inches(0), SW, SH, fill=BG_TEAL)

# Accent line
_add_rect(s, Inches(0.8), Inches(1.0), Inches(8.4), Inches(0.06), fill=TEAL)

# Title
_add_text_box(
    s,
    Inches(0.8),
    Inches(1.3),
    Inches(8.4),
    Inches(1.5),
    [
        ("MindPulse", Pt(40), True, WHITE),
        (
            "Privacy-First Behavioral Stress Detection System",
            Pt(20),
            False,
            RGBColor(0xAA, 0xDD, 0xEE),
        ),
    ],
)

# Subtitle
_add_text_box(
    s,
    Inches(0.8),
    Inches(3.2),
    Inches(8.4),
    Inches(0.6),
    [
        (
            'Detecting Stress from "HOW You Type", Not "WHAT You Type"',
            Pt(16),
            False,
            RGBColor(0x88, 0xBB, 0xCC),
        ),
    ],
)

# Divider
_add_rect(s, Inches(0.8), Inches(4.1), Inches(2.0), Inches(0.04), fill=WHITE)

# Team info card
_add_card(
    s,
    Inches(0.8),
    Inches(4.4),
    Inches(8.4),
    Inches(2.2),
    [
        ("Team Members", Pt(13), True, TEAL),
        (
            "Neha (01111502723)  |  Anand Misra (04011502723)  |  Pratham Nahar (05311502723)",
            Pt(12),
            False,
            DARK,
        ),
        ("", Pt(6)),
        ("Mentored by: Dr. Jolly Parikh", Pt(12), True, DARK),
        ("Mini Project Review — Academic Year 2024-25", Pt(11), False, MID_GRAY),
    ],
    bg=WHITE,
    border_color=TEAL,
)

_set_notes(
    s,
    "Welcome to our Mini Project Review on MindPulse. Our system detects workplace stress from typing patterns, mouse movements, and app-switching behavior without ever capturing what the user types.",
)

# ================================================================
# SLIDE 2: INTRODUCTION / PROBLEM
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[1])
_clear_slide(s)

# Title bar
_add_rect(s, Inches(0), Inches(0), SW, Inches(1.0), fill=TEAL)
_add_text_box(
    s,
    Inches(0.5),
    Inches(0.15),
    Inches(9.0),
    Inches(0.7),
    [
        ("Introduction / Problem Statement", Pt(28), True, WHITE),
    ],
)

# Three cards
problems = [
    (
        "The Crisis",
        [
            ("83% of US workers experience", Pt(11), False, DARK),
            ("workplace stress", Pt(11), True, RED),
            ("", Pt(4)),
            ("$300 billion annual cost to", Pt(11), False, DARK),
            ("employers (AIS, 2023)", Pt(11), True, RED),
        ],
    ),
    (
        "Current Methods",
        [
            ("Self-reports: subjective,", Pt(11), False, DARK),
            ("infrequent, recall bias", Pt(10), False, MID_GRAY),
            ("", Pt(4)),
            ("Wearables: invasive,", Pt(11), False, DARK),
            ("expensive hardware", Pt(10), False, MID_GRAY),
            ("", Pt(4)),
            ("Physiological: impractical", Pt(11), False, DARK),
            ("for continuous use", Pt(10), False, MID_GRAY),
        ],
    ),
    (
        "Research Gap",
        [
            ("No existing system detects", Pt(11), False, DARK),
            ("stress from computer", Pt(11), False, DARK),
            ("interaction while", Pt(11), False, DARK),
            ("preserving privacy", Pt(13), True, TEAL),
        ],
    ),
]
for i, (title, lines) in enumerate(problems):
    x = Inches(0.5 + i * 3.2)
    all_lines = [(title, Pt(14), True, TEAL), ("", Pt(4))] + lines
    _add_card(s, x, Inches(1.3), Inches(2.9), Inches(2.5), all_lines, border_color=TEAL)

# Bottom insight box
_add_card(
    s,
    Inches(0.5),
    Inches(4.2),
    Inches(9.0),
    Inches(2.5),
    [
        ("Our Key Insight", Pt(16), True, TEAL),
        ("", Pt(6)),
        (
            'MindPulse captures "HOW you interact" — not "WHAT you type"',
            Pt(14),
            False,
            DARK,
        ),
        ("", Pt(4)),
        ("We analyze behavioral metadata only:", Pt(12), True, DARK),
        ("  Keyboard: timing patterns, hold/flight times, rhythm", Pt(11), False, DARK),
        (
            "  Mouse: movement speed, click patterns, scroll behavior",
            Pt(11),
            False,
            DARK,
        ),
        ("  Context: app-switch frequency, session fragmentation", Pt(11), False, DARK),
    ],
    border_color=TEAL,
)

_set_notes(
    s,
    "Workplace stress is a $300B problem. Current detection methods fail: self-reports are subjective, wearables are invasive. MindPulse addresses the gap by detecting stress from behavioral metadata while guaranteeing privacy.",
)

# ================================================================
# SLIDE 3: SOLUTION OVERVIEW
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[1])
_clear_slide(s)

_add_rect(s, Inches(0), Inches(0), SW, Inches(1.0), fill=TEAL)
_add_text_box(
    s,
    Inches(0.5),
    Inches(0.15),
    Inches(9.0),
    Inches(0.7),
    [
        ("Our Solution: MindPulse", Pt(28), True, WHITE),
    ],
)

# Feature cards in 2 rows
features = [
    ("3-Class Prediction", "NEUTRAL / MILD / STRESSED\nstress level classification"),
    (
        "Zero Content Capture",
        "Never records typed characters\nonly metadata and timing",
    ),
    (
        "On-Device Processing",
        "All inference happens locally\nno data leaves your machine",
    ),
    ("Personal Adaptation", "Learns individual baseline\nvia EMA + SQLite calibration"),
    ("Real-Time Monitoring", "WebSocket streaming with\n5-minute rolling windows"),
    (
        "23 Feature Vector",
        "Most papers use 8-12; we use 23\nincluding 3 novel features",
    ),
]
for i, (title, desc) in enumerate(features):
    col = i % 3
    row = i // 3
    x = Inches(0.5 + col * 3.2)
    y = Inches(1.3 + row * 2.8)
    _add_card(
        s,
        x,
        y,
        Inches(2.9),
        Inches(2.5),
        [
            (title, Pt(14), True, TEAL),
            ("", Pt(4)),
            (desc, Pt(11), False, DARK),
        ],
        border_color=TEAL,
    )

_set_notes(
    s,
    "MindPulse has six key capabilities: 3-class stress prediction, zero content capture, on-device processing, personal adaptation, real-time monitoring, and a 23-feature behavioral vector with 3 novel features.",
)

# ================================================================
# SLIDE 4: OBJECTIVES
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[1])
_clear_slide(s)

_add_rect(s, Inches(0), Inches(0), SW, Inches(1.0), fill=TEAL)
_add_text_box(
    s,
    Inches(0.5),
    Inches(0.15),
    Inches(9.0),
    Inches(0.7),
    [
        ("Project Objectives", Pt(28), True, WHITE),
    ],
)

objectives = [
    ("1", "Design a privacy-preserving stress detection system"),
    ("2", "Analyze user behavioral patterns (keyboard, mouse, context switching)"),
    (
        "3",
        "Extract key stress-related features (typing irregularity, error rate, rage clicks)",
    ),
    ("4", "Develop a hybrid machine learning model for stress classification"),
    ("5", "Enable real-time, on-device stress monitoring"),
    ("6", "Implement personalized calibration with adaptive stress interventions"),
]

for i, (num, text) in enumerate(objectives):
    y = Inches(1.3 + i * 0.95)
    # Number circle
    circ = s.shapes.add_shape(
        MSO_SHAPE.OVAL, Inches(0.5), y + Inches(0.05), Inches(0.5), Inches(0.5)
    )
    circ.fill.solid()
    circ.fill.fore_color.rgb = TEAL
    circ.line.fill.background()
    circ.shadow.inherit = False
    tf = circ.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    _set_para(run, num, size=Pt(16), bold=True, color=WHITE)

    # Objective text
    _add_text_box(
        s,
        Inches(1.2),
        y,
        Inches(6.5),
        Inches(0.5),
        [
            (text, Pt(14), False, DARK),
        ],
    )

    # Status badge
    badge = _add_rounded(
        s, Inches(7.9), y + Inches(0.05), Inches(1.6), Inches(0.4), fill=TEAL
    )
    tf2 = badge.text_frame
    tf2.paragraphs[0].alignment = PP_ALIGN.CENTER
    run2 = tf2.paragraphs[0].add_run()
    _set_para(run2, "Completed", Pt(10), bold=True, color=WHITE)

_set_notes(
    s,
    "All six original objectives completed. Privacy designed, behavioral patterns analyzed, 23 features extracted, XGBoost model trained, real-time monitoring via WebSocket, personal calibration with EMA.",
)

# ================================================================
# SLIDE 5: ROADMAP
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[1])
_clear_slide(s)

_add_rect(s, Inches(0), Inches(0), SW, Inches(1.0), fill=TEAL)
_add_text_box(
    s,
    Inches(0.5),
    Inches(0.15),
    Inches(9.0),
    Inches(0.7),
    [
        ("Project Roadmap", Pt(28), True, WHITE),
    ],
)

phases = [
    (
        "Phase 1",
        "Research & Design",
        "Literature review, behavioral indicators, privacy architecture",
        True,
    ),
    (
        "Phase 2",
        "Data & Features",
        "pynput collector, 23-feature pipeline, synthetic data generator",
        True,
    ),
    (
        "Phase 3",
        "Model Development",
        "XGBoost, DualNormalizer, PersonalBaseline, 1D-CNN",
        True,
    ),
    (
        "Phase 4",
        "Application",
        "FastAPI backend, Next.js frontend, WebSocket streaming",
        True,
    ),
    ("Phase 5", "Evaluation", "Group K-Fold, calibration, benchmark comparison", True),
    (
        "Phase 6",
        "Real Deployment",
        "15-20 real users, browser extension, multi-modal fusion",
        False,
    ),
]

for i, (phase, title, desc, done) in enumerate(phases):
    y = Inches(1.2 + i * 1.0)
    # Phase label
    color = TEAL if done else RGBColor(0xE0, 0xE0, 0xE0)
    txt_color = WHITE if done else DARK
    phase_box = _add_rounded(s, Inches(0.3), y, Inches(1.4), Inches(0.8), fill=color)
    tf = phase_box.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    run = tf.paragraphs[0].add_run()
    _set_para(run, phase, size=Pt(12), bold=True, color=txt_color)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    _set_para(run2, "Done" if done else "Next", size=Pt(9), bold=False, color=txt_color)

    # Connector line
    if i < len(phases) - 1:
        _add_rect(
            s,
            Inches(0.95),
            y + Inches(0.8),
            Inches(0.04),
            Inches(0.2),
            fill=TEAL if done else LIGHT,
        )

    # Title
    _add_text_box(
        s,
        Inches(1.9),
        y,
        Inches(2.5),
        Inches(0.4),
        [
            (title, Pt(14), True, DARK),
        ],
    )

    # Description
    _add_text_box(
        s,
        Inches(1.9),
        y + Inches(0.35),
        Inches(7.5),
        Inches(0.4),
        [
            (desc, Pt(11), False, MID_GRAY),
        ],
    )

_set_notes(
    s,
    "Six phases. Phases 1-5 complete. Phase 6 (real-user deployment) is next: 15-20 volunteers for 2-4 weeks each.",
)

# ================================================================
# SLIDE 6: ARCHITECTURE
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[1])
_clear_slide(s)

_add_rect(s, Inches(0), Inches(0), SW, Inches(1.0), fill=TEAL)
_add_text_box(
    s,
    Inches(0.5),
    Inches(0.15),
    Inches(9.0),
    Inches(0.7),
    [
        ("System Architecture", Pt(28), True, WHITE),
    ],
)

# Four main boxes
arch_boxes = [
    (
        Inches(0.3),
        Inches(1.5),
        "Desktop App",
        "Python/pynput\n\npynput collector\nKeyboard/Mouse\nContext events",
    ),
    (
        Inches(2.8),
        Inches(1.5),
        "ML Engine",
        "XGBoost + CNN\n\n23-feature extraction\nDualNormalizer\n46-dim input",
    ),
    (
        Inches(5.3),
        Inches(1.5),
        "Backend API",
        "FastAPI\n\n7 REST endpoints\nWebSocket\nPort 5000",
    ),
    (
        Inches(5.3),
        Inches(4.2),
        "Frontend",
        "Next.js 15\n\nRiskMeter gauge\nTimelineChart\nMetrics cards",
    ),
]

for left, top, title, content in arch_boxes:
    card = _add_rounded(s, left, top, Inches(2.3), Inches(2.3), fill=BG_TEAL)
    tf = card.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.15)
    p = tf.paragraphs[0]
    run = p.add_run()
    _set_para(run, title, size=Pt(16), bold=True, color=WHITE)
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    run2 = p2.add_run()
    _set_para(run2, content, size=Pt(10), bold=False, color=RGBColor(0xAA, 0xDD, 0xEE))

# Arrows
arrow_positions = [
    (Inches(2.6), Inches(2.6), Inches(0.2), Inches(0.08)),
    (Inches(5.1), Inches(2.6), Inches(0.2), Inches(0.08)),
    (Inches(6.4), Inches(3.8), Inches(0.08), Inches(0.4)),
]
for l, t, w, h in arrow_positions:
    _add_rect(s, l, t, w, h, fill=TEAL)

_set_notes(
    s,
    "Four components: Desktop App captures events, ML Engine extracts features and runs XGBoost, FastAPI backend serves predictions, Next.js frontend provides real-time visualization.",
)

# ================================================================
# SLIDE 7: DATA COLLECTION
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[1])
_clear_slide(s)

_add_rect(s, Inches(0), Inches(0), SW, Inches(1.0), fill=TEAL)
_add_text_box(
    s,
    Inches(0.5),
    Inches(0.15),
    Inches(9.0),
    Inches(0.7),
    [
        ("Data Collection (Privacy-First)", Pt(28), True, WHITE),
    ],
)

# Left: What we record
_add_card(
    s,
    Inches(0.3),
    Inches(1.3),
    Inches(4.5),
    Inches(5.5),
    [
        ("What We Record", Pt(16), True, TEAL),
        ("", Pt(6)),
        ("Keyboard:", Pt(13), True, DARK),
        ("  timestamp_press, timestamp_release", Pt(11), False, DARK),
        ("  key_category (alpha/digit/special)", Pt(11), False, DARK),
        ("", Pt(4)),
        ("Mouse:", Pt(13), True, DARK),
        ("  x, y position, timestamp", Pt(11), False, DARK),
        ("  click_type, scroll_delta", Pt(11), False, DARK),
        ("", Pt(4)),
        ("Context:", Pt(13), True, DARK),
        ("  switch_timestamp", Pt(11), False, DARK),
        ("  app_category_hash (SHA-256)", Pt(11), False, DARK),
    ],
    border_color=TEAL,
)

# Right: What we DON'T record
_add_card(
    s,
    Inches(5.2),
    Inches(1.3),
    Inches(4.5),
    Inches(2.8),
    [
        ("What We DON'T Record", Pt(16), True, RED),
        ("", Pt(6)),
        ("  Actual key characters", Pt(12), False, DARK),
        ("  Screen content or URLs", Pt(12), False, DARK),
        ("  Application names (hashed)", Pt(12), False, DARK),
        ("  File contents or messages", Pt(12), False, DARK),
    ],
    border_color=RED,
)

_add_card(
    s,
    Inches(5.2),
    Inches(4.4),
    Inches(4.5),
    Inches(2.4),
    [
        ("Privacy Guarantees", Pt(16), True, TEAL),
        ("", Pt(6)),
        ("  Zero content capture", Pt(12), False, DARK),
        ("  On-device processing", Pt(12), False, DARK),
        ("  SHA-256 app hashing", Pt(12), False, DARK),
        ("  No data leaves machine", Pt(12), False, DARK),
    ],
    border_color=TEAL,
)

_set_notes(
    s,
    "Privacy is our core design. We capture only timing metadata. App names are SHA-256 hashed. Zero content capture ensures complete privacy.",
)

# ================================================================
# SLIDE 8: FEATURES
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[1])
_clear_slide(s)

_add_rect(s, Inches(0), Inches(0), SW, Inches(1.0), fill=TEAL)
_add_text_box(
    s,
    Inches(0.5),
    Inches(0.15),
    Inches(9.0),
    Inches(0.7),
    [
        ("Feature Extraction (23 Dimensions)", Pt(28), True, WHITE),
    ],
)

categories = [
    (
        "Keyboard (11)",
        "hold_time_mean/std/median\nflight_time_mean/std\ntyping_speed_wpm\nerror_rate\npause_frequency/duration\nburst_length_mean\nrhythm_entropy",
    ),
    (
        "Mouse (6)",
        "mouse_speed_mean/std\ndirection_change_rate\nclick_count\nrage_click_count\nscroll_velocity_std",
    ),
    ("Context (3)", "tab_switch_freq\nswitch_entropy\nsession_fragmentation"),
    ("Temporal (3)", "hour_of_day\nday_of_week\nsession_duration_min"),
]

for i, (title, content) in enumerate(categories):
    x = Inches(0.3 + i * 2.45)
    _add_card(
        s,
        x,
        Inches(1.3),
        Inches(2.25),
        Inches(4.0),
        [
            (title, Pt(14), True, TEAL),
            ("", Pt(4)),
            (content, Pt(10), False, DARK),
        ],
        border_color=TEAL,
    )

# Bottom bar
_add_card(
    s,
    Inches(0.3),
    Inches(5.6),
    Inches(9.4),
    Inches(1.2),
    [
        ("Dual Normalization Pipeline", Pt(14), True, TEAL),
        (
            "23 raw features  ->  [23 global z-scores] + [23 per-user circadian z-scores]  ->  46 features",
            Pt(11),
            False,
            DARK,
        ),
    ],
    bg=LIGHT,
    border_color=TEAL,
)

_set_notes(
    s,
    "23 features across 4 categories. Dual normalization creates 46-dimensional input by concatenating global and per-user z-scores.",
)

# ================================================================
# SLIDE 9: MODEL
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[1])
_clear_slide(s)

_add_rect(s, Inches(0), Inches(0), SW, Inches(1.0), fill=TEAL)
_add_text_box(
    s,
    Inches(0.5),
    Inches(0.15),
    Inches(9.0),
    Inches(0.7),
    [
        ("Machine Learning Model & Calibration", Pt(28), True, WHITE),
    ],
)

# Left: Model
_add_card(
    s,
    Inches(0.3),
    Inches(1.3),
    Inches(4.5),
    Inches(5.5),
    [
        ("XGBoost Classifier", Pt(16), True, TEAL),
        ("", Pt(4)),
        ("Objective: multi:softprob (3-class)", Pt(11), False, DARK),
        ("Max depth: 6", Pt(11), False, DARK),
        ("Estimators: 350", Pt(11), False, DARK),
        ("Learning rate: 0.08", Pt(11), False, DARK),
        ("Subsample: 0.9", Pt(11), False, DARK),
        ("Class weights: Balanced", Pt(11), False, DARK),
        ("", Pt(6)),
        ("1D-CNN (Secondary Branch)", Pt(14), True, TEAL),
        ("", Pt(4)),
        ("Processes raw keystroke sequences", Pt(11), False, DARK),
        ("Captures digraph/trigraph patterns", Pt(11), False, DARK),
    ],
    border_color=TEAL,
)

# Right: Calibration
_add_card(
    s,
    Inches(5.2),
    Inches(1.3),
    Inches(4.5),
    Inches(2.6),
    [
        ("DualNormalizer", Pt(14), True, TEAL),
        ("", Pt(4)),
        ("Global: (x - mu) / sigma", Pt(11), False, DARK),
        ("Per-user circadian hourly baselines", Pt(11), False, DARK),
        ("Output: 46-dim feature vector", Pt(11), False, DARK),
    ],
    border_color=TEAL,
)

_add_card(
    s,
    Inches(5.2),
    Inches(4.2),
    Inches(4.5),
    Inches(2.6),
    [
        ("PersonalBaseline (SQLite + EMA)", Pt(14), True, TEAL),
        ("", Pt(4)),
        ("Per-user feature means per hour", Pt(11), False, DARK),
        ("Exponential Moving Average updates", Pt(11), False, DARK),
        ("+1.9% F1 improvement", Pt(13), True, TEAL),
    ],
    border_color=TEAL,
)

_set_notes(
    s,
    "XGBoost with 350 estimators. DualNormalizer creates 46 features. PersonalBaseline uses SQLite with EMA for per-user calibration (+1.9% F1).",
)

# ================================================================
# SLIDE 10: LITERATURE REVIEW
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[1])
_clear_slide(s)

_add_rect(s, Inches(0), Inches(0), SW, Inches(1.0), fill=TEAL)
_add_text_box(
    s,
    Inches(0.5),
    Inches(0.15),
    Inches(9.0),
    Inches(0.7),
    [
        ("Literature Review", Pt(28), True, WHITE),
    ],
)

headers = ["Year & Citation", "Technology", "Dataset", "Key Finding", "Result"]
data = [
    [
        "Naegelin et al.\n(2025) ETH Zurich",
        "Keyboard+mouse\nML analysis",
        "OSF: qpekf\n36 emp, 8wk",
        "Universal models\nperform poorly",
        "F1 = 7.8%",
    ],
    [
        "VTT Finland\n(2024)",
        "AI mouse\nmovement",
        "Proprietary\nworkplace",
        "Mouse > heart\nrate for stress",
        "71% acc",
    ],
    [
        "Pepa et al.\n(2021)",
        "Keystroke\ndynamics RF",
        "Zenodo\n62 users",
        "In-the-wild\nfeasible",
        "F1 = 60%",
    ],
    [
        "CMU (2023)",
        "Hold time +\nlatency",
        "CMU InfSci\n116 subjects",
        "Fits-and-starts\npattern",
        "76% acc",
    ],
    [
        "ETH Zurich\n(2023)",
        "Neuromotor\nnoise theory",
        "OSF: qpekf\nLab study",
        "Typing > heart\nrate",
        "F1 = 62.5%",
    ],
]

_add_table(
    s,
    Inches(0.3),
    Inches(1.3),
    Inches(9.4),
    Inches(4.5),
    headers,
    data,
    col_widths=[Inches(1.8), Inches(1.6), Inches(1.8), Inches(2.2), Inches(1.5)],
)

_set_notes(
    s,
    "Literature review: ETH Zurich 2025 found universal models perform poorly. VTT Finland showed mouse alone achieves 71%. Pepa et al. demonstrated in-the-wild feasibility with 60% F1.",
)

# ================================================================
# SLIDE 11: PROGRESS
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[1])
_clear_slide(s)

_add_rect(s, Inches(0), Inches(0), SW, Inches(1.0), fill=TEAL)
_add_text_box(
    s,
    Inches(0.5),
    Inches(0.15),
    Inches(9.0),
    Inches(0.7),
    [
        ("Progress Since Last Review", Pt(28), True, WHITE),
    ],
)

headers = ["Component", "Status", "Details"]
data = [
    ["ML Core", "Complete", "4 Python modules: collector, extractor, model, synthetic"],
    ["Backend API", "Complete", "FastAPI: 7 REST endpoints + WebSocket streaming"],
    ["Frontend", "Complete", "Next.js 15: 6 pages, 7 components, real-time WS"],
    ["Evaluation", "Complete", "Group K-Fold, calibration eval, realistic simulator"],
    ["Documentation", "Complete", "Research report, ML pipeline, results summary"],
    ["GitHub", "Live", "53 files at iAMv1/mindpulse, clean structure"],
]

_add_table(
    s,
    Inches(0.3),
    Inches(1.3),
    Inches(9.4),
    Inches(4.5),
    headers,
    data,
    col_widths=[Inches(1.8), Inches(1.2), Inches(6.4)],
)

_set_notes(
    s,
    "All components completed since last review: ML core, backend API, frontend dashboard, evaluation pipeline, documentation, and GitHub repository.",
)

# ================================================================
# SLIDE 12: CHALLENGES
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[1])
_clear_slide(s)

_add_rect(s, Inches(0), Inches(0), SW, Inches(1.0), fill=TEAL)
_add_text_box(
    s,
    Inches(0.5),
    Inches(0.15),
    Inches(9.0),
    Inches(0.7),
    [
        ("Challenges Faced & Overcome", Pt(28), True, WHITE),
    ],
)

challenges = [
    (
        "Feature Mismatch",
        "Model expected 46, inference passed 23",
        "Implemented DualNormalizer",
    ),
    (
        "Import Failures",
        "Relative imports broke in __main__",
        "Converted to package imports",
    ),
    (
        "Duplicate Servers",
        "Two separate FastAPI implementations",
        "Consolidated into single backend",
    ),
    (
        "Unrealistic Metrics",
        "99%+ F1 from separable classes",
        "Created overlap simulator",
    ),
    ("Unicode Errors", "Box-drawing chars caused codec errors", "Replaced with ASCII"),
    (
        "Path Resolution",
        "Artifacts pointed to wrong directory",
        "Module-relative paths",
    ),
]

# Column headers
_add_text_box(
    s,
    Inches(0.3),
    Inches(1.15),
    Inches(2.5),
    Inches(0.4),
    [("Challenge", Pt(12), True, TEAL)],
)
_add_text_box(
    s,
    Inches(3.0),
    Inches(1.15),
    Inches(3.2),
    Inches(0.4),
    [("Problem", Pt(12), True, RED)],
)
_add_text_box(
    s,
    Inches(6.5),
    Inches(1.15),
    Inches(3.0),
    Inches(0.4),
    [("Solution", Pt(12), True, TEAL)],
)

for i, (ch, prob, sol) in enumerate(challenges):
    y = Inches(1.6 + i * 0.85)
    _add_card(
        s,
        Inches(0.3),
        y,
        Inches(2.5),
        Inches(0.7),
        [
            (ch, Pt(11), True, DARK),
        ],
        border_color=TEAL,
    )
    _add_card(
        s,
        Inches(3.0),
        y,
        Inches(3.2),
        Inches(0.7),
        [
            (prob, Pt(10), False, DARK),
        ],
        border_color=RED,
    )
    _add_card(
        s,
        Inches(6.5),
        y,
        Inches(3.0),
        Inches(0.7),
        [
            (sol, Pt(10), False, TEAL),
        ],
        border_color=TEAL,
    )

_set_notes(
    s,
    "Six major challenges: feature shape mismatch, import failures, duplicate servers, unrealistic metrics, Unicode errors, path resolution. All successfully resolved.",
)

# ================================================================
# SLIDE 13: METRICS
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[1])
_clear_slide(s)

_add_rect(s, Inches(0), Inches(0), SW, Inches(1.0), fill=TEAL)
_add_text_box(
    s,
    Inches(0.5),
    Inches(0.15),
    Inches(9.0),
    Inches(0.7),
    [
        ("Performance Comparison & Metrics", Pt(28), True, WHITE),
    ],
)

# Results table
headers = ["Evaluation", "F1-Macro", "Accuracy", "Std"]
data = [
    ["Universal (no cal)", "46.8%", "48.0%", "+/- 4.6%"],
    ["Calibrated (40 samp)", "48.7%", "50.0%", "+/- 4.0%"],
    ["Improvement", "+1.9%", "+2.0%", "--"],
]
_add_table(s, Inches(0.3), Inches(1.3), Inches(4.5), Inches(1.8), headers, data)

# Per-class
_add_card(
    s,
    Inches(0.3),
    Inches(3.4),
    Inches(4.5),
    Inches(2.3),
    [
        ("Per-Class Performance", Pt(13), True, TEAL),
        ("", Pt(4)),
        ("NEUTRAL   P=0.54  R=0.71  F1=0.62", Pt(10), False, DARK),
        ("MILD      P=0.31  R=0.11  F1=0.17", Pt(10), False, RED),
        ("STRESSED  P=0.28  R=0.33  F1=0.30", Pt(10), False, DARK),
    ],
    border_color=TEAL,
)

# Benchmark comparison
headers2 = ["Study", "F1-Macro"]
data2 = [
    ["ETH Zurich 2025 (Universal)", "7.8%"],
    ["MindPulse (Universal)", "46.8%"],
    ["MindPulse (Calibrated)", "48.7%"],
    ["ETH Zurich 2023 (Lab)", "62.5%"],
    ["Pepa et al. 2021", "60.0%"],
]
_add_table(s, Inches(5.2), Inches(1.3), Inches(4.5), Inches(3.0), headers2, data2)

_add_card(
    s,
    Inches(5.2),
    Inches(4.6),
    Inches(4.5),
    Inches(1.6),
    [
        ("Key Observation", Pt(13), True, RED),
        ("", Pt(4)),
        ("MILD class bottleneck (F1=0.17)", Pt(11), False, DARK),
        ("Heavy overlap with NEUTRAL/STRESSED", Pt(10), False, MID_GRAY),
    ],
    border_color=RED,
)

_set_notes(
    s,
    "Universal model: 46.8% F1 matching ETH Zurich 2025. Calibrated: 48.7% (+1.9%). MILD class is bottleneck at F1=0.17.",
)

# ================================================================
# SLIDE 14: GAPS
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[1])
_clear_slide(s)

_add_rect(s, Inches(0), Inches(0), SW, Inches(1.0), fill=TEAL)
_add_text_box(
    s,
    Inches(0.5),
    Inches(0.15),
    Inches(9.0),
    Inches(0.7),
    [
        ("Research Gaps", Pt(28), True, WHITE),
    ],
)

# Left: Overcome
_add_card(
    s,
    Inches(0.3),
    Inches(1.3),
    Inches(4.5),
    Inches(5.5),
    [
        ("Gaps We Overcame", Pt(16), True, TEAL),
        ("", Pt(6)),
        ("  Limited feature sets (8-12) -> 23 features", Pt(11), False, DARK),
        ("  No personal calibration -> DualNormalizer + EMA", Pt(11), False, DARK),
        ("  Privacy concerns -> Zero content capture", Pt(11), False, DARK),
        ("  Simplified evaluation -> Group K-Fold", Pt(11), False, DARK),
        ("  No self-report noise -> 76.8% accuracy sim", Pt(11), False, DARK),
        ("  No real-time capability -> WebSocket streaming", Pt(11), False, DARK),
    ],
    border_color=TEAL,
)

# Right: New
_add_card(
    s,
    Inches(5.2),
    Inches(1.3),
    Inches(4.5),
    Inches(5.5),
    [
        ("New Research Gaps", Pt(16), True, RED),
        ("", Pt(6)),
        ("  No real-user data (all simulated)", Pt(11), False, DARK),
        ("  MILD class detection weak (F1=0.17)", Pt(11), False, DARK),
        ("  Desktop-only (pynput dependency)", Pt(11), False, DARK),
        ("  Calibration cold start problem", Pt(11), False, DARK),
        ("  Long-term concept drift untested", Pt(11), False, DARK),
    ],
    border_color=RED,
)

_set_notes(
    s,
    "Overcame 6 gaps: limited features, no calibration, privacy, evaluation, noise modeling, real-time. 5 new gaps: real-user data, MILD class, platform, cold start, drift.",
)

# ================================================================
# SLIDE 15: NEW CHALLENGES
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[1])
_clear_slide(s)

_add_rect(s, Inches(0), Inches(0), SW, Inches(1.0), fill=TEAL)
_add_text_box(
    s,
    Inches(0.5),
    Inches(0.15),
    Inches(9.0),
    Inches(0.7),
    [
        ("New Challenges & Strategy", Pt(28), True, WHITE),
    ],
)

headers = ["Challenge", "Problem", "Strategy", "Timeline"]
data = [
    [
        "Real-User Data",
        "All eval on simulated",
        "Deploy 15-20 volunteers",
        "2-3 months",
    ],
    ["MILD Class", "F1=0.17, overlap", "Ordinal regression", "1 month"],
    ["Cross-Platform", "Desktop-only", "Chrome extension", "2 months"],
    ["Cold Start", "No baseline", "Pop -> personal model", "Done"],
    ["Concept Drift", "Behavior changes", "SPC + retraining", "3 months"],
]

_add_table(
    s,
    Inches(0.3),
    Inches(1.3),
    Inches(9.4),
    Inches(3.5),
    headers,
    data,
    col_widths=[Inches(1.8), Inches(2.0), Inches(3.0), Inches(1.6)],
)

_set_notes(
    s,
    "Five challenges with strategies: real-user deployment, MILD ordinal regression, Chrome extension, population-to-personal, concept drift retraining.",
)

# ================================================================
# SLIDE 16: CONCLUSION
# ================================================================
s = prs.slides.add_slide(prs.slide_layouts[1])
_clear_slide(s)

_add_rect(s, Inches(0), Inches(0), SW, Inches(1.0), fill=TEAL)
_add_text_box(
    s,
    Inches(0.5),
    Inches(0.15),
    Inches(9.0),
    Inches(0.7),
    [
        ("Conclusion & References", Pt(28), True, WHITE),
    ],
)

# Left: Takeaways
_add_card(
    s,
    Inches(0.3),
    Inches(1.3),
    Inches(4.5),
    Inches(3.0),
    [
        ("Key Takeaways", Pt(16), True, TEAL),
        ("", Pt(4)),
        ("  Privacy-first stress detection is feasible", Pt(12), False, DARK),
        ("  46.8% F1 universal, 48.7% calibrated", Pt(12), False, DARK),
        ("  23 features with 3 novel ones", Pt(12), False, DARK),
        ("  Complete full-stack implementation", Pt(12), False, DARK),
        ("  Honest evaluation (Group K-Fold)", Pt(12), False, DARK),
    ],
    border_color=TEAL,
)

_add_card(
    s,
    Inches(0.3),
    Inches(4.6),
    Inches(4.5),
    Inches(2.2),
    [
        ("Next Steps", Pt(14), True, TEAL),
        ("", Pt(4)),
        ("  Real-user deployment (15-20 users)", Pt(11), False, DARK),
        ("  Browser extension development", Pt(11), False, DARK),
        ("  Multi-modal fusion (facial + physio)", Pt(11), False, DARK),
    ],
    border_color=TEAL,
)

# Right: References
refs = [
    "[1] Naegelin et al. (2025) ETH Zurich - osf.io/qpekf",
    "[2] VTT Finland (2024) AI mouse movement",
    "[3] Pepa et al. (2021) Keystroke dynamics - Zenodo",
    "[4] CMU (2023) Keystroke stress - CMU InfSci",
    "[5] ETH Zurich (2023) Neuromotor noise - OSF",
    "[6] XGBoost - github.com/dmlc/xgboost",
    "[7] FastAPI - github.com/tiangolo/fastapi",
    "[8] Next.js - github.com/vercel/next.js",
    "[9] pynput - pypi.org/project/pynput",
    "[10] scikit-learn - github.com/scikit-learn",
    "[11] Amer. Inst. Stress (2023) - stress.org",
]
_add_card(
    s,
    Inches(5.2),
    Inches(1.3),
    Inches(4.5),
    Inches(5.5),
    [
        ("References", Pt(14), True, TEAL),
        ("", Pt(4)),
        *[(ref, Pt(9), False, DARK) for ref in refs],
    ],
    border_color=TEAL,
)

_set_notes(
    s,
    "MindPulse demonstrates privacy-first stress detection is feasible. 46.8% F1 universal, 48.7% calibrated. Complete full-stack with honest evaluation. Next: real deployment, browser extension, multi-modal.",
)

# ================================================================
# SAVE
# ================================================================
output = "D:/another-project/MindPulse_Review.pptx"
prs.save(output)
print(f"Saved: {output}")
print(f"Slides: {len(prs.slides)}")
